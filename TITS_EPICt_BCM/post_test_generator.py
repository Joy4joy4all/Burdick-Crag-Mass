# -*- coding: utf-8 -*-
#!/usr/bin/env python3
"""
POST-TEST GENERATOR — Per-Test Run Slide + Report Engine
================================================================
Generates after EVERY completed test_run:

Output (per test) into {project}/test_report_slides/:
  1. Test Run_{N}_{Person}_Log_Slide.pptx   — BCM format log slide
  2. Test Run_{N}_{Person}_Compound_Report.txt — 12-section delta report

MODULAR: Zero hardcoded project IDs. Uses project_paths for all routing.
EXCEL-CONNECTED: Pulls hypotheses from the test_run plan Excel.

© 2025-2026 Stephen J. Burdick Sr. — All Rights Reserved.
For all the industrial workers lost to preventable acts.
"""

import json, re, sys
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional
from collections import defaultdict

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# ── Import project_paths if available ──
try:
    from project_paths import (
        get_project_dir, get_report_slide_dir, get_tested_dir,
        find_project_config, list_deployed_projects, find_project_baseline_deck,
    )
    PROJECT_PATHS_AVAILABLE = True
except ImportError:
    PROJECT_PATHS_AVAILABLE = False
    def get_project_dir(pid): return Path.cwd() / "BCM_Projects" / "BCM_Navigator_Project" / pid
    def get_report_slide_dir(pid): return get_project_dir(pid) / "test_report_slides"
    def get_tested_dir(pid): return get_project_dir(pid) / "tested"
    def find_project_config(pid): return None
    def find_project_baseline_deck(pid): return None
    def list_deployed_projects(): return []


# ══════════════════════════════════════════════════════════════
# COLORS — BCM purple format (matches uploaded slides exactly)
# ══════════════════════════════════════════════════════════════
C_PURPLE_DARK = "3B0764"
C_PURPLE_MED  = "7B2D8E"
C_LAVENDER    = "D5C6E0"
C_WHITE = "FFFFFF"; C_BLACK = "000000"
C_RED = "C41E3A"; C_GREEN = "2E7D32"; C_GRAY = "666666"
C_RESULTS_BG = RGBColor(0xD4, 0xED, 0xEC)
C_ACTION_BG  = RGBColor(0xD9, 0xEA, 0xD3)


# ══════════════════════════════════════════════════════════════
# VP + BMC KEYWORD MAPS
# ══════════════════════════════════════════════════════════════
VP_DEFINITIONS = {
    "VP1_DAMAGE_PREVENTION":     {"short": "Damage Prevention",      "keywords": ["damage","rock","destroy","metal","tramp","catastrophic","broken","worn","replacement","repair","failure","knife","knives","wipe out"]},
    "VP2_PLANNED_MAINTENANCE":   {"short": "Planned Maintenance",    "keywords": ["planned","preventive","predictive","schedule","reactive","emergency","unplanned","readings"]},
    "VP3_OPERATOR_VISIBILITY":   {"short": "Operator Visibility",    "keywords": ["visibility","real-time","sensor","monitor","detect","blind","no data","flying blind","can't see","no visibility"]},
    "VP4_CHEMICAL_OPTIMIZATION": {"short": "Chemical Optimization",  "keywords": ["chemical","ClO2","bleach","cooking","kappa","yield","consistency","optimization","savings"]},
    "VP5_SUPPLIER_ACCOUNTABILITY":{"short":"Supplier Accountability","keywords": ["vendor","supplier","accountability","log","trace","claim","premium","certified","verification","dock"]},
    "VP6_PITCH_CONTAMINATION":   {"short": "Pitch Contamination",    "keywords": ["pitch","felt","paper machine","press","wire","deposit","stickies","fourdrinier","nip"]},
}

BMC_BLOCKS = {
    "Key Partnerships":       ["university","partner","OEM","vendor","supplier","research","alliance","institute"],
    "Key Activities":         ["maintenance","knife","replacement","screen","inspection","cleaning","repair","monitoring","install","calibrat"],
    "Key Resources":          ["sensor","algorithm","AI","acoustic","patent","data","calibration","technology"],
    "Value Propositions":     ["damage","prevent","detect","save","protect","visibility","real-time","automated","uptime"],
    "Customer Relationships": ["trust","data","dashboard","report","subscription","service","support","demo"],
    "Channels":               ["direct","outreach","pilot","trade show","referral","network","contact","regional"],
    "Substrate Classs":      ["mill","kraft","operator","manager","superintendent","foreman","buyer","producer","chip","pulp"],
    "Cost Structure":         ["hardware","installation","sensor cost","hosting","development","certification","software"],
    "Revenue Streams":        ["subscription","license","SaaS","per-unit","contract","annual","ROI","payback","buy"],
}


# ══════════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════════
def _count_kw(text, kws):
    t = (text or "").lower()
    return sum(1 for k in kws if k.lower() in t)

def _safe_trunc(text, n=500):
    s = str(text or "").strip()
    return s[:n] + "..." if len(s) > n else s

def _parse_cost(raw):
    if not raw: return 0.0
    try: return float(raw)
    except: pass
    nums = re.findall(r'[\d,]+', str(raw).replace('$', ''))
    if nums:
        try: return float(nums[0].replace(',', ''))
        except: pass
    return 0.0

def _all_text(p):
    parts = [p.get('results', ''), p.get('action_iterate', ''),
             p.get('hypotheses_validation', ''), p.get('hypotheses', '')]
    for q, a in p.get('experiments_answers', {}).items():
        parts.extend([str(q), str(a)])
    for eq in p.get('equipment_discussed', []):
        parts.append(eq.get('notes', ''))
    return ' '.join(str(x) for x in parts if x)


# ══════════════════════════════════════════════════════════════
# EXCEL HYPOTHESIS CONNECTOR
# ══════════════════════════════════════════════════════════════
def load_hypotheses_from_excel(project_id, person_name=None, test_num=None):
    """
    Pull hypotheses from project's Test Run Plan Excel.
    Matches by person name (primary), test_num as fallback for legacy data.
    """
    excel_path = find_project_config(project_id)
    if not excel_path or not Path(excel_path).exists():
        return {}
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(excel_path), read_only=True, data_only=True)
        ws = wb.active
        best_match = {}
        for r in range(2, ws.max_row + 1):
            row_num = ws.cell(r, 1).value
            row_name = str(ws.cell(r, 2).value or "")
            row_title = str(ws.cell(r, 3).value or "")
            row_type = str(ws.cell(r, 4).value or "")
            row_company = str(ws.cell(r, 5).value or "")
            row_hyp = str(ws.cell(r, 6).value or "")
            row_q = str(ws.cell(r, 7).value or "")

            if person_name and row_name:
                pn = person_name.lower().strip()
                rn = row_name.lower().strip()
                if pn in rn or rn in pn:
                    wb.close()
                    return {'hypothesis_context': row_hyp, 'planned_questions': row_q,
                            'excel_name': row_name, 'excel_title': row_title,
                            'excel_type': row_type, 'excel_company': row_company,
                            'match_row': r}
                elif len(pn.split()) > 0 and len(rn.split()) > 0:
                    if pn.split()[-1] in rn or rn.split()[-1] in pn:
                        if not best_match:
                            best_match = {'hypothesis_context': row_hyp, 'planned_questions': row_q,
                                          'excel_name': row_name, 'excel_title': row_title,
                                          'excel_type': row_type, 'excel_company': row_company,
                                          'match_row': r}

            # Legacy fallback: match by test_num
            if not best_match and test_num is not None and row_num == test_num:
                best_match = {'hypothesis_context': row_hyp, 'planned_questions': row_q,
                              'excel_name': row_name, 'excel_title': row_title,
                              'excel_type': row_type, 'excel_company': row_company,
                              'match_row': r}

        wb.close()
        return best_match
    except Exception as e:
        print(f"  [WARN] Excel hypothesis load: {e}")
        return {}


# ══════════════════════════════════════════════════════════════
# PPTX HELPERS
# ══════════════════════════════════════════════════════════════
def _gradient_header(slide, prs):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.65))
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(0x3B, 0x07, 0x64)
    bar.line.fill.background()
    bar2 = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.22))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = RGBColor(0x5B, 0x21, 0x7F)
    bar2.line.fill.background()

def _rich_tb(slide, l, t, w, h, runs, sz=11, align=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align if align is not None else PP_ALIGN.LEFT
    for text, bold, uline in runs:
        r = p.add_run(); r.text = str(text)
        r.font.size = Pt(sz); r.font.name = "Calibri"
        r.font.color.rgb = RGBColor(0, 0, 0)
        r.font.bold = bold; r.font.underline = uline
    return tb

def _section_tb(slide, l, t, w, label, content, sz=11, max_c=450):
    text = _safe_trunc(content, max_c) or "[Not filled]"
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(2.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = f"{label} "
    r1.font.size = Pt(sz); r1.font.name = "Calibri"; r1.font.bold = True
    r1.font.color.rgb = RGBColor(0, 0, 0)
    r2 = p.add_run(); r2.text = text
    r2.font.size = Pt(sz); r2.font.name = "Calibri"; r2.font.bold = False
    r2.font.color.rgb = RGBColor(0, 0, 0)
    nlines = max(1, len(text) // 85 + 1)
    return t + nlines * 0.22 + 0.15

def _simple_tb(slide, l, t, w, h, text, sz=11, bold=False, chex="000000", fn="Calibri", al=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = str(text)
    p.font.size = Pt(sz); p.font.bold = bold; p.font.name = fn
    p.font.color.rgb = RGBColor.from_string(chex)
    p.alignment = al if al is not None else PP_ALIGN.LEFT
    return tb


# ══════════════════════════════════════════════════════════════
# MAIN CLASS
# ══════════════════════════════════════════════════════════════
class PostTest RunGenerator:
    def __init__(self, database, project_id: str):
        self.database = database
        self.project_id = project_id
        self.output_dir = get_report_slide_dir(project_id)
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def generate_all(self, parsed, intel_before=None, intel_after=None):
        person = parsed.get('script_name', 'Unknown')
        pfx = person.replace(' ', '_')

        excel_hyp = load_hypotheses_from_excel(
            self.project_id, person_name=person,
            test_num=parsed.get('test_num', 0))
        if excel_hyp:
            parsed['_excel_hypothesis'] = excel_hyp
            print(f"  [EXCEL] Matched row {excel_hyp.get('match_row')} for {person}")

        outputs = {}
        tasks = [
            ('log_slide',       self._gen_log_slide),
            ('compound_report', lambda p, o: self._gen_report(p, intel_before, intel_after, o)),
        ]
        for name, method in tasks:
            ext = '.txt' if 'report' in name else '.pptx'
            suffix = {'log_slide': 'Log_Slide', 'compound_report': 'Compound_Report'}[name]
            path = self.output_dir / f"{pfx}_{suffix}{ext}"
            try:
                method(parsed, path)
                outputs[name] = path
                print(f"  [OK] {path.name}")
            except Exception as e:
                print(f"  [ERR] {name}: {e}")
                import traceback; traceback.print_exc()
        print(f"[PostTest RunGen] {len(outputs)}/2 for #{num} -> {self.output_dir}")
        return outputs

    # ════════════════════════════════════════════════════════════
    # SLIDE 1: BCM TEST LOG (purple format)
    # ════════════════════════════════════════════════════════════
    def _gen_log_slide(self, pd, out):
        if not PPTX_AVAILABLE:
            raise RuntimeError("python-pptx not installed")

        prs = Presentation()
        prs.slide_width = Inches(10); prs.slide_height = Inches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(0xD5, 0xC6, 0xE0)
        _gradient_header(slide, prs)

        hdr = pd.get('header', {})
        num = pd.get('test_num', 0)
        person = hdr.get('Person', pd.get('script_name', ''))
        title = hdr.get('Title', pd.get('title', ''))
        company = hdr.get('Company', pd.get('source_version', ''))
        ctype = hdr.get('Substrate Class', pd.get('test_category', ''))
        date_s = hdr.get('Date', pd.get('date', datetime.now().strftime('%m/%d/%y')))

        _rich_tb(slide, 0.5, 0.75, 3.0, 0.4, [("Test Run Log", True, True)], sz=14)

        box = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(9.0), Inches(1.35))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
        box.line.color.rgb = RGBColor(0, 0, 0); box.line.width = Pt(1.0)

        y = 1.25
        _rich_tb(slide, 0.6, y, 9.0, 0.3, [
            ("Team Name: ", False, False), ("GIBUSH", False, True),
            ("     Test Run Date:   ", False, False), (str(date_s), False, True),
            ("          Test Run Number : ", False, False), (str(num), False, True)])
        y += 0.28
        _rich_tb(slide, 0.6, y, 9.0, 0.3, [
            ("Person Test Runed (Name): ", False, False), (str(person), False, True),
            ("   Title (Position): ", False, False), (str(title), False, True)])
        y += 0.28
        _rich_tb(slide, 0.6, y, 9.0, 0.3, [
            ("Company: ", False, False), (str(company), False, True)])
        y += 0.28
        _rich_tb(slide, 0.6, y, 9.0, 0.3, [
            ("Customer Type (User, Decision Maker, Payer, etc.): ", False, False),
            (str(ctype), False, True)])

        y = 2.75
        excel_hyp = pd.get('_excel_hypothesis', {})
        hyp = excel_hyp.get('hypothesis_context', '') or \
              pd.get('hypotheses_validation', '') or pd.get('hypotheses', '')
        y = _section_tb(slide, 0.6, y, 8.8, "Hypotheses:", hyp)
        y += 0.15

        exps = pd.get('experiments_answers', {})
        exp_text = pd.get('experiments', '') or \
                   '. '.join(f"Asked: {q}" for q in list(exps.keys())[:3]) or '[Not filled]'
        y = _section_tb(slide, 0.6, y, 8.8, "Experiments:", exp_text)
        y += 0.15

        results = pd.get('results', '')
        if results and "[FILL" not in results:
            bg = slide.shapes.add_shape(1, Inches(0.3), Inches(y-0.05), Inches(9.4), Inches(1.2))
            bg.fill.solid(); bg.fill.fore_color.rgb = C_RESULTS_BG; bg.line.fill.background()
        y = _section_tb(slide, 0.6, y, 8.8, "Results:", results)
        y += 0.15

        action = pd.get('action_iterate', '')
        if action and "[FILL" not in action:
            bg = slide.shapes.add_shape(1, Inches(0.3), Inches(y-0.05), Inches(9.4), Inches(0.9))
            bg.fill.solid(); bg.fill.fore_color.rgb = C_ACTION_BG; bg.line.fill.background()
        _section_tb(slide, 0.6, y, 8.8, "Action/Iterate:", action)

        prs.save(str(out))

    # ════════════════════════════════════════════════════════════
    # 12-SECTION COMPOUND REPORT
    # ════════════════════════════════════════════════════════════
    def _gen_report(self, pd, intel_b, intel_a, out):
        num = pd.get('test_num', 0)
        person = pd.get('script_name', '?')
        company = pd.get('source_version', '?')
        cube = pd.get('header', {}).get('Q-Cube Position', '')
        if not cube:
            ql = pd.get('q_layer', ''); qo = pd.get('q_object', '')
            qs = pd.get('q_stack', [])
            if ql: cube = f"[{ql}, {qo}, {','.join(str(s) for s in qs) if isinstance(qs, list) else qs}]"

        results = pd.get('results', '')
        action = pd.get('action_iterate', '')
        hyp_val = pd.get('hypotheses_validation', '') or pd.get('hypotheses', '')
        at = _all_text(pd)
        equip = pd.get('equipment_discussed', [])
        equip_impacts = pd.get('substrate_impacts', [])
        costs = pd.get('equipment_costs', [])
        exps = pd.get('experiments_answers', {})
        pains = pd.get('pain_indicators_noted', [])
        has_i = bool(intel_b and intel_a)
        excel_hyp = pd.get('_excel_hypothesis', {})

        L = []
        L.append("=" * 80)
        L.append(f"COMPOUND INTELLIGENCE REPORT — Test Run #{num}: {person}")
        L.append(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        L.append(f"Company: {company}"); L.append(f"Cube Position: {cube}")
        if excel_hyp: L.append(f"Excel Match: Row {excel_hyp.get('match_row', '?')}")
        L.append("=" * 80)

        # S1
        L.extend(["", "─" * 80, "SECTION 1: TEST SUMMARY", "─" * 80])
        L.append(f"\nResults:\n{_safe_trunc(results, 600)}")
        L.append(f"\nHypothesis Validation:\n{_safe_trunc(hyp_val, 400)}")
        L.append(f"\nAction/Iterate:\n{_safe_trunc(action, 400)}")

        # S2
        L.extend(["", "─" * 80, "SECTION 2: EQUIPMENT DISCUSSED + DAMAGE DATA", "─" * 80])
        if equip:
            for eq in equip:
                n = f" -> {eq['notes']}" if eq.get('notes') else ''
                L.append(f"  * {eq['equipment']} ({eq.get('typical_cost', 'N/A')}){n}")
        elif equip_impacts:
            for eq in equip_impacts:
                if isinstance(eq, dict):
                    L.append(f"  * {eq.get('equipment', eq.get('name', '?'))}")
        if not equip and not equip_impacts: L.append("  No equipment data.")

        # S3
        L.extend(["", "─" * 80, "SECTION 3: EXPERIMENT ANSWERS", "─" * 80])
        if exps:
            for q, a in exps.items():
                L.append(f"\n  Q: {_safe_trunc(q, 120)}")
                L.append(f"  A: {_safe_trunc(a, 200) or '(no answer)'}")
        else: L.append("  None recorded.")

        # S4
        L.extend(["", "─" * 80, "SECTION 4: COMPLETENESS SCORE", "─" * 80])
        checks = [
            ("Results filled",        bool(results and "[FILL" not in results)),
            ("Equipment checklist",   len(equip) > 0 or len(equip_impacts) > 0),
            ("Hypothesis validation", bool(hyp_val)),
            ("Action items",          bool(action and "[FILL" not in action)),
            ("Experiment answers",    len(exps) > 0),
            ("Cost data",             len(costs) > 0),
            ("Pain indicators",       len(pains) > 0),
        ]
        score = sum(1 for _, p in checks if p)
        for lb, p in checks: L.append(f"  {'[OK]' if p else '[  ]'} {lb}")
        L.append(f"\n  TOTAL: {score}/{len(checks)} ({score / len(checks) * 100:.0f}%)")

        # S5
        L.extend(["", "─" * 80, "SECTION 5: VALUE PROPOSITION DELTA", "─" * 80])
        L.append("  (Keyword scan)")
        for vk, vd in VP_DEFINITIONS.items():
            h = _count_kw(at, vd['keywords'])
            tag = " << STRONG" if h >= 3 else " << evidence" if h >= 1 else ""
            L.append(f"  {vd['short']:28s} {h:2d} keyword hits{tag}")

        # S6
        L.extend(["", "─" * 80, "SECTION 6: HYPOTHESIS DELTA", "─" * 80])
        if excel_hyp:
            L.append(f"  [FROM EXCEL] {_safe_trunc(excel_hyp.get('hypothesis_context', ''), 300)}")
        L.append(f"  {_safe_trunc(hyp_val, 400)}")

        # S7
        L.extend(["", "─" * 80, "SECTION 7: EQUIPMENT DAMAGE RANKING DELTA", "─" * 80])
        if equip:
            for eq in equip: L.append(f"  * {eq['equipment']}: {eq.get('notes', '')}")
        else: L.append("  No equipment ranking.")

        # S8
        L.extend(["", "─" * 80, "SECTION 8: CUBE DENSITY", "─" * 80])
        if self.database and hasattr(self.database, 'tests'):
            lc = defaultdict(int)
            for iv in self.database.test_runs:
                ly = getattr(iv, 'q_layer', '')
                if ly: lc[ly] += 1
            for l, c in sorted(lc.items()): L.append(f"  {l}: {c} test_runs")
        L.append(f"  This test_run: {cube}")

        # S9
        L.extend(["", "─" * 80, "SECTION 9: BMC IMPACT", "─" * 80])
        for bn, kws in BMC_BLOCKS.items():
            h = _count_kw(at, kws)
            if h >= 3:
                best = max(kws, key=lambda k: at.lower().count(k.lower()))
                L.append(f"  [+++] {bn}: {h} hits ('{best}')")
            elif h >= 1: L.append(f"  [+]   {bn}: {h} hit(s)")
        empty = [b for b, k in BMC_BLOCKS.items() if _count_kw(at, k) == 0]
        if empty: L.append(f"\n  No evidence: {', '.join(empty)}")

        # S10
        L.extend(["", "─" * 80, "SECTION 10: CUSTOMER SEGMENT UPDATE", "─" * 80])
        if self.database and hasattr(self.database, 'tests'):
            lc = defaultdict(int); lco = defaultdict(set)
            for iv in self.database.test_runs:
                ly = getattr(iv, 'q_layer', '')
                if ly: lc[ly] += 1; lco[ly].add(getattr(iv, 'source_version', ''))
            tot = sum(lc.values())
            for l in sorted(lc):
                cos = ', '.join(sorted(c for c in lco[l] if c)[:4])
                L.append(f"  {l}: {lc[l]} ({lc[l] / max(tot, 1) * 100:.0f}%) — {cos}")
        bh = _count_kw(at, ["budget", "amo", "capex", "roi", "payback"])
        if bh > 0: L.append(f"\n  Budget PATH signals: {bh}")
        else: L.append(f"\n  No budget signals — probe in follow-up")

        # S11
        L.extend(["", "─" * 80, "SECTION 11: CONTRADICTIONS", "─" * 80])
        L.append("  No computed contradictions.")
        if _count_kw(at, ["no one cares", "don't care"]) > 0 and \
           _count_kw(at, ["operators sense", "they know"]) > 0:
            L.append("  !! POSSIBLE: Apathy vs awareness — probe next test_run")

        # S12
        L.extend(["", "─" * 80, "SECTION 12: NEXT TEST RECOMMENDATIONS", "─" * 80])
        if not any(c.get('annual_damage') for c in costs):
            L.append(f"  Dollar gap: $0 from {person} — FOLLOW UP")
        L.append(f"  Completeness: {score}/{len(checks)}")
        wb = [b for b, k in BMC_BLOCKS.items() if _count_kw(at, k) == 0]
        if wb: L.append(f"  BMC gaps: {', '.join(wb[:3])}")

        L.extend(["", "=" * 80, "END OF COMPOUND INTELLIGENCE REPORT", "=" * 80])
        out.parent.mkdir(parents=True, exist_ok=True)
        out.write_text('\n'.join(L), encoding='utf-8')


# ══════════════════════════════════════════════════════════════
# BATCH GENERATOR
# ══════════════════════════════════════════════════════════════
def batch_generate_from_deck(project_id, deck_path, database=None):
    """Generate slides for all tests in a Baseline deck JSON."""
    with open(deck_path, 'r', encoding='utf-8') as f:
        deck = json.load(f)
    tests = deck.get('tests', deck if isinstance(deck, list) else [])
    if not test_runs:
        print(f"No tests in {deck_path}"); return 0

    gen = PostTest RunGenerator(database, project_id)
    total = 0
    for iv in test_runs:
        results = iv.get('results', '')
        if not results or '[PENDING' in results:
            print(f"  SKIP {iv.get('script_name', '?')}: no results"); continue

        parsed = {
            'test_num': iv.get('test_num', 0),
            'script_name': iv.get('script_name', ''), 'source_version': iv.get('source_version', ''),
            'title': iv.get('title', ''), 'test_category': iv.get('test_category', ''),
            'date': iv.get('date', ''),
            'header': {
                'Person': iv.get('script_name', ''), 'Title': iv.get('title', ''),
                'Company': iv.get('source_version', ''),
                'Substrate Class': iv.get('test_category', ''),
                'Date': iv.get('date', ''),
                'Q-Cube Position': f"[{iv.get('q_layer','')}, {iv.get('q_object','')}, {iv.get('q_stack','')}]",
            },
            'results': results, 'hypotheses': iv.get('hypotheses', ''),
            'hypotheses_validation': iv.get('hypotheses', ''),
            'experiments': iv.get('experiments', ''),
            'action_iterate': iv.get('action_iterate', ''),
            'experiments_answers': {}, 'equipment_discussed': [],
            'substrate_impacts': iv.get('substrate_impacts', []),
            'equipment_costs': [], 'pain_indicators_noted': [],
            'q_layer': iv.get('q_layer', ''), 'q_object': iv.get('q_object', ''),
            'q_stack': iv.get('q_stack', []),
        }
        outputs = gen.generate_all(parsed)
        if outputs: total += 1

    print(f"\n[BATCH] {total}/{len(tests)} slides generated for {project_id}")
    return total


if __name__ == "__main__":
    print("=" * 60)
    print("POST-TEST GENERATOR — Standalone Test")
    print("=" * 60)
    projects = list_deployed_projects()
    if not projects: projects = ["BCM_SUBSTRATE"]
    for pid in projects:
        print(f"\n[PROJECT] {pid}")
        deck_path = find_project_baseline_deck(pid)
        if deck_path:
            print(f"  Baseline deck: {deck_path}")
            batch_generate_from_deck(pid, deck_path)
        else:
            print(f"  No Baseline deck found")
